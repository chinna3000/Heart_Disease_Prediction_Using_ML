from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

PPTX_PATH = r"Heart_Disease_Prediction.pptx"

prs = Presentation(PPTX_PATH)

# Print existing slides
print("Existing slides:")
for i, slide in enumerate(prs.slides):
    title = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape == slide.shapes.title:
            title = shape.text_frame.text
            break
        elif shape.has_text_frame:
            title = shape.text_frame.paragraphs[0].text
            break
    print(f"  Slide {i+1}: {title}")

# Find the methodology slide index
methodology_idx = None
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if "methodology" in para.text.lower():
                    methodology_idx = i
                    break
        if methodology_idx is not None:
            break
    if methodology_idx is not None:
        break

print(f"\nMethodology slide found at index: {methodology_idx}")

# ── We'll insert the new slide AFTER the methodology slide ──────────────────
# python-pptx doesn't support insert, so we append then move.

slide_layout = prs.slide_layouts[1]  # Title and Content layout

new_slide = prs.slides.add_slide(slide_layout)

# ── BACKGROUND ────────────────────────────────────────────────────────────────
background = new_slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy

# ── TITLE ─────────────────────────────────────────────────────────────────────
title_box = new_slide.shapes.title
title_box.text = "Algorithm: Random Forest Classifier"
tf = title_box.text_frame
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)   # gold
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ── CONTENT BOX (overview + steps) ───────────────────────────────────────────
# Remove existing placeholder content box and recreate for full control
for shape in new_slide.placeholders:
    if shape.placeholder_format.idx == 1:
        sp = shape._element
        sp.getparent().remove(sp)
        break

slide_w = prs.slide_width
slide_h = prs.slide_height

# Overview section
overview_box = new_slide.shapes.add_textbox(
    Inches(0.4), Inches(1.35), slide_w - Inches(0.8), Inches(0.65)
)
tf_ov = overview_box.text_frame
tf_ov.word_wrap = True

p = tf_ov.paragraphs[0]
p.text = "Random Forest is an ensemble learning method that builds multiple decision trees during training and outputs the class that is the mode of the classes (classification) of the individual trees."
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0xD0, 0xE8, 0xFF)
p.alignment = PP_ALIGN.LEFT

# Steps / Algorithm box
steps_box = new_slide.shapes.add_textbox(
    Inches(0.4), Inches(2.15), slide_w - Inches(0.8), Inches(4.2)
)
tf_steps = steps_box.text_frame
tf_steps.word_wrap = True

algorithm_steps = [
    ("Step 1 – Data Collection & Preprocessing",
     "Load the Heart Disease dataset (270 records, 8 features). Handle missing values, encode categorical variables (Heart Disease: Presence / Absence)."),
    ("Step 2 – Feature Selection",
     "Features used: Age, Chest Pain Type, Blood Pressure, Cholesterol, Max HR, ST Depression, Number of Vessels Fluro, Thallium."),
    ("Step 3 – Dataset Splitting",
     "Split data into Training Set (80%) and Testing Set (20%) using train_test_split."),
    ("Step 4 – Build Decision Trees (Bootstrap Sampling)",
     "Randomly sample subsets of training data with replacement (bagging). For each subset, grow a Decision Tree using a random subset of features at each node split (feature randomness)."),
    ("Step 5 – Voting / Aggregation",
     "Each tree independently predicts the class label. Final prediction = majority vote across all N trees."),
    ("Step 6 – Model Evaluation",
     "Evaluate on test set: Accuracy ≈ 98.5%, Precision, Recall, F1-Score computed via classification report."),
    ("Step 7 – Model Persistence & Deployment",
     "Save trained model using Pickle → heartdiseaseprediction.model. Load model in Flask web app for real-time prediction."),
]

first = True
for heading, detail in algorithm_steps:
    if first:
        p_head = tf_steps.paragraphs[0]
        first = False
    else:
        p_head = tf_steps.add_paragraph()

    p_head.text = heading
    p_head.font.bold = True
    p_head.font.size = Pt(12)
    p_head.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)   # gold headings
    p_head.space_before = Pt(6)

    p_det = tf_steps.add_paragraph()
    p_det.text = "  " + detail
    p_det.font.size = Pt(11)
    p_det.font.color.rgb = RGBColor(0xD0, 0xE8, 0xFF)
    p_det.space_after = Pt(2)

# ── MOVE the new slide to just after the methodology slide ────────────────────
# python-pptx stores slides in prs.slides._sldIdLst
xml_slides = prs.slides._sldIdLst
# The new slide is currently the last element
new_slide_elem = xml_slides[-1]
xml_slides.remove(new_slide_elem)

target_position = (methodology_idx + 1) if methodology_idx is not None else len(xml_slides)
xml_slides.insert(target_position, new_slide_elem)

prs.save(PPTX_PATH)
print(f"\n✅ Algorithm slide added successfully after slide {methodology_idx + 1} (Methodology).")
print(f"   PPTX saved to: {PPTX_PATH}")
